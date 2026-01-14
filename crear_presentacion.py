from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- DISEÑO ---
    # Función auxiliar para títulos
    def add_title(slide, text):
        title = slide.shapes.title
        title.text = text
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80) # Azul oscuro

    # Función auxiliar para contenido
    def add_content(slide, text_list):
        content_box = slide.placeholders[1]
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for item in text_list:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(24)
            p.space_after = Pt(14)

    # --- SLIDE 1: TÍTULO ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "WEEKENDER AI"
    subtitle.text = "Tu Planificador de Viajes Inteligente\nProyecto Full-Stack Data Engineering & GenAI"
    
    # --- SLIDE 2: EL PROBLEMA ---
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    add_title(slide, "El Problema: 'Parálisis por Análisis'")
    content = [
        "• Planificar un fin de semana conlleva horas de búsqueda.",
        "• Demasiada información dispersa (vuelos, hoteles, blogs).",
        "• Dificultad para ajustar planes a un presupuesto fijo.",
        "• SOLUCIÓN: Un asistente que personaliza itinerarios en segundos."
    ]
    add_content(slide, content)

    # --- SLIDE 3: ARQUITECTURA TÉCNICA (STACK) ---
    slide = prs.slides.add_slide(slide_layout)
    add_title(slide, "Arquitectura Tecnológica")
    content = [
        "FRONTEND (La Interfaz):",
        "   - Python Streamlit + Custom CSS.",
        "   - Gestión de estado (Session State) para el chat.",
        "",
        "BACKEND (El Motor):",
        "   - Flask API (RESTful Service).",
        "   - Orquestación de lógica y seguridad."
    ]
    add_content(slide, content)

    # --- SLIDE 4: INTELIGENCIA & DATOS ---
    slide = prs.slides.add_slide(slide_layout)
    add_title(slide, "IA Generativa y Persistencia")
    content = [
        "INTELIGENCIA ARTIFICIAL (LLM):",
        "   - API de Groq (Modelo Llama-3.3-70b).",
        "   - Inferencia de ultra-baja latencia.",
        "   - Prompt Engineering (System & User prompts).",
        "",
        "BASE DE DATOS:",
        "   - PostgreSQL (Alojada en Render Cloud).",
        "   - Almacenamiento de logs e historial de usuarios."
    ]
    add_content(slide, content)

    # --- SLIDE 5: FLUJO DE DATOS (PIPELINE) ---
    slide = prs.slides.add_slide(slide_layout)
    add_title(slide, "Flujo de Datos (Pipeline)")
    content = [
        "1. INGESTA: Usuario envía destino y presupuesto (Streamlit).",
        "2. PROCESADO: Flask recibe JSON y construye el prompt.",
        "3. INFERENCIA: Groq genera el itinerario estructurado.",
        "4. PERSISTENCIA: Se guarda la interacción en PostgreSQL.",
        "5. RESPUESTA: El itinerario se muestra en el Frontend."
    ]
    add_content(slide, content)

    # --- SLIDE 6: INFRAESTRUCTURA Y DEPLOY ---
    slide = prs.slides.add_slide(slide_layout)
    add_title(slide, "Infraestructura & MLOps")
    content = [
        "• DOCKERIZACIÓN:",
        "   - Contenedor único con Python 3.9-slim.",
        "   - Script de arranque dual (Flask + Streamlit).",
        "",
        "• GESTIÓN DE ENTORNO:",
        "   - Variables de entorno (.env) para seguridad de API Keys.",
        "",
        "• TESTING:",
        "   - Tests automatizados de endpoints (Pytest)."
    ]
    add_content(slide, content)

    # --- SLIDE 7: DEMO ---
    slide_layout = prs.slide_layouts[0] # Title Only essentially
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "LIVE DEMO"
    subtitle = slide.placeholders[1]
    subtitle.text = "Demostración del sistema en tiempo real"

    # --- GUARDAR ---
    prs.save('Presentacion_Weekender.pptx')
    print("✅ ¡Presentación creada con éxito! Abre el archivo 'Presentacion_Weekender.pptx'")

if __name__ == "__main__":
    create_presentation()